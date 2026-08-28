import pytest

import fitz
from pptx import Presentation
from pptx.util import Inches

from app.services.extraction import ExtractionError, extract_pdf, extract_pptx, extract_txt


def test_extract_txt_reads_and_cleans_whitespace(tmp_path):
    path = tmp_path / "notes.txt"
    path.write_text("  Line one.  \n\n\n  Line two.  \n", encoding="utf-8")

    segments = extract_txt(path, source_name="notes.txt")

    assert len(segments) == 1
    seg = segments[0]
    assert seg.source_type == "txt"
    assert seg.source_name == "notes.txt"
    assert seg.page is None
    assert seg.slide is None
    assert "Line one." in seg.text
    assert "Line two." in seg.text


def test_extract_txt_raises_on_empty_file(tmp_path):
    path = tmp_path / "empty.txt"
    path.write_text("   \n\n   ", encoding="utf-8")

    with pytest.raises(ExtractionError, match="empty"):
        extract_txt(path, source_name="empty.txt")


def test_extract_txt_raises_on_invalid_utf8(tmp_path):
    path = tmp_path / "bad_encoding.txt"
    path.write_bytes(b"\xff\xfe\x00\x01invalid")

    with pytest.raises(ExtractionError, match="UTF-8"):
        extract_txt(path, source_name="bad_encoding.txt")


def _make_pdf(tmp_path, pages_text: list[str]):
    doc = fitz.open()
    for text in pages_text:
        page = doc.new_page()
        if text:
            page.insert_text((72, 72), text)
    path = tmp_path / "sample.pdf"
    doc.save(str(path))
    doc.close()
    return path


def test_extract_pdf_returns_one_segment_per_page(tmp_path):
    path = _make_pdf(tmp_path, ["First page body", "Second page body"])

    segments = extract_pdf(path, source_name="sample.pdf")

    assert len(segments) == 2
    assert segments[0].page == 1
    assert segments[1].page == 2
    assert "First page body" in segments[0].text
    assert "Second page body" in segments[1].text
    assert all(seg.source_type == "pdf" for seg in segments)
    assert all(seg.source_name == "sample.pdf" for seg in segments)


def test_extract_pdf_skips_blank_pages_but_keeps_real_page_numbers(tmp_path):
    path = _make_pdf(tmp_path, ["Only real content", ""])

    segments = extract_pdf(path, source_name="sample.pdf")

    assert len(segments) == 1
    assert segments[0].page == 1


def test_extract_pdf_raises_when_no_extractable_text(tmp_path):
    path = _make_pdf(tmp_path, ["", ""])

    with pytest.raises(ExtractionError, match="no extractable text"):
        extract_pdf(path, source_name="sample.pdf")


def test_extract_pdf_raises_on_corrupted_file(tmp_path):
    path = tmp_path / "corrupted.pdf"
    path.write_bytes(b"this is not a real pdf")

    with pytest.raises(ExtractionError, match="Corrupted"):
        extract_pdf(path, source_name="corrupted.pdf")


def _make_pptx(tmp_path, slides: list[dict]):
    """slides: list of {"heading": str|None, "table_rows": list[list[str]]|None}"""
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]

    for slide_spec in slides:
        slide = prs.slides.add_slide(blank_layout)

        heading = slide_spec.get("heading")
        if heading:
            textbox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(6), Inches(1))
            textbox.text_frame.text = heading

        table_rows = slide_spec.get("table_rows")
        if table_rows:
            rows, cols = len(table_rows), len(table_rows[0])
            table_shape = slide.shapes.add_table(rows, cols, Inches(1), Inches(2), Inches(6), Inches(1.5))
            table = table_shape.table
            for r, row_values in enumerate(table_rows):
                for c, value in enumerate(row_values):
                    table.cell(r, c).text = value

    path = tmp_path / "sample.pptx"
    prs.save(str(path))
    return path


def test_extract_pptx_returns_one_segment_per_slide_with_text(tmp_path):
    path = _make_pptx(tmp_path, [
        {"heading": "Slide one heading"},
        {"heading": "Slide two heading"},
    ])

    segments = extract_pptx(path, source_name="sample.pptx")

    assert len(segments) == 2
    assert segments[0].slide == 1
    assert segments[1].slide == 2
    assert "Slide one heading" in segments[0].text
    assert "Slide two heading" in segments[1].text
    assert all(seg.source_type == "pptx" for seg in segments)


def test_extract_pptx_includes_table_cell_text(tmp_path):
    path = _make_pptx(tmp_path, [
        {
            "heading": "Pricing",
            "table_rows": [["Plan", "Price"], ["Basic", "$10"], ["Pro", "$25"]],
        },
    ])

    segments = extract_pptx(path, source_name="sample.pptx")

    assert len(segments) == 1
    text = segments[0].text
    assert "Pricing" in text
    assert "Basic" in text and "$10" in text
    assert "Pro" in text and "$25" in text


def test_extract_pptx_raises_when_no_extractable_text(tmp_path):
    path = _make_pptx(tmp_path, [{}, {}])

    with pytest.raises(ExtractionError, match="no extractable text"):
        extract_pptx(path, source_name="sample.pptx")


def test_extract_pptx_raises_on_corrupted_file(tmp_path):
    path = tmp_path / "corrupted.pptx"
    path.write_bytes(b"this is not a real pptx")

    with pytest.raises(ExtractionError, match="Corrupted"):
        extract_pptx(path, source_name="corrupted.pptx")
