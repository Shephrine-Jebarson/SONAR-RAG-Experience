"""Tests for checkpoint 3: embeddings service, vector_store service,
/upload, /add-url, /process routes.

All Gemini and Qdrant calls are mocked — no real API keys required.
"""

import io
import time
import uuid
from unittest.mock import MagicMock, patch

import fitz
import pytest
from fastapi.testclient import TestClient
from pptx import Presentation
from pptx.util import Inches

from app.dependencies import get_gemini_client, get_httpx_client, get_qdrant_client
from app.main import app
from app.models import Chunk, ExtractedSegment
from app.services.embeddings import EMBEDDING_DIM, EmbeddingError, embed_chunks, embed_query
from app.services.vector_store import VectorStoreError, ensure_collection, upsert_chunks


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _fake_gemini_client(dim: int = EMBEDDING_DIM):
    """Return a mock genai.Client whose embed_content returns dim-length vectors."""
    client = MagicMock()

    def _embed(model, contents, config=None):
        response = MagicMock()
        response.embeddings = [
            MagicMock(values=[0.1] * dim) for _ in contents
        ]
        return response

    client.models.embed_content.side_effect = _embed
    return client


def _fake_qdrant_client():
    client = MagicMock()
    client.get_collections.return_value = MagicMock(collections=[])
    return client


def _make_pdf_bytes(text: str) -> bytes:
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 72), text)
    buf = io.BytesIO()
    doc.save(buf)
    doc.close()
    return buf.getvalue()


def _make_pptx_bytes(heading: str) -> bytes:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
    tb.text_frame.text = heading
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# Embedding service tests
# ---------------------------------------------------------------------------

def test_embed_chunks_returns_one_vector_per_chunk():
    chunks = [
        Chunk(chunk_id="a_000_000", text="hello", source_type="txt", source_name="a.txt"),
        Chunk(chunk_id="b_000_000", text="world", source_type="txt", source_name="b.txt"),
    ]
    client = _fake_gemini_client()

    vectors = embed_chunks(chunks, client)

    assert len(vectors) == 2
    assert all(len(v) == EMBEDDING_DIM for v in vectors)


def test_embed_chunks_empty_input_returns_empty():
    client = _fake_gemini_client()
    assert embed_chunks([], client) == []


def test_embed_chunks_raises_embedding_error_on_api_failure():
    chunks = [Chunk(chunk_id="x_000_000", text="text", source_type="txt", source_name="x.txt")]
    client = MagicMock()
    client.models.embed_content.side_effect = RuntimeError("quota exceeded")

    with pytest.raises(EmbeddingError, match="quota exceeded"):
        embed_chunks(chunks, client)


def test_embed_query_returns_single_vector():
    client = _fake_gemini_client()
    vector = embed_query("what is RAG?", client)
    assert len(vector) == EMBEDDING_DIM


# ---------------------------------------------------------------------------
# Vector store tests
# ---------------------------------------------------------------------------

def test_ensure_collection_creates_when_missing():
    qdrant = _fake_qdrant_client()
    ensure_collection(qdrant, "test_col")
    qdrant.create_collection.assert_called_once()


def test_ensure_collection_tolerates_concurrent_create_race():
    """Two requests hitting /reset or /process near-simultaneously (e.g. React
    StrictMode double-invoking an effect, or two browser tabs) can both see
    "collection missing" and both call create_collection — Qdrant accepts the
    first and 409s the second. ensure_collection must treat that 409 as
    success (the desired state — the collection exists — was still achieved)
    instead of surfacing it as a VectorStoreError."""
    from qdrant_client.http.exceptions import UnexpectedResponse

    qdrant = MagicMock()
    qdrant.get_collections.return_value = MagicMock(collections=[])  # not yet created
    qdrant.create_collection.side_effect = UnexpectedResponse(
        status_code=409,
        reason_phrase="Conflict",
        content=b'{"status":{"error":"Wrong input: Collection `test_col` already exists!"}}',
        headers={},
    )

    ensure_collection(qdrant, "test_col")  # must not raise


def test_ensure_collection_reraises_non_conflict_errors():
    from qdrant_client.http.exceptions import UnexpectedResponse

    qdrant = MagicMock()
    qdrant.get_collections.return_value = MagicMock(collections=[])
    qdrant.create_collection.side_effect = UnexpectedResponse(
        status_code=500,
        reason_phrase="Internal Server Error",
        content=b"boom",
        headers={},
    )

    with pytest.raises(VectorStoreError):
        ensure_collection(qdrant, "test_col")


def test_ensure_collection_skips_when_already_exists():
    qdrant = MagicMock()
    existing = MagicMock()
    existing.name = "test_col"
    qdrant.get_collections.return_value = MagicMock(collections=[existing])

    ensure_collection(qdrant, "test_col")

    qdrant.create_collection.assert_not_called()


def test_upsert_chunks_calls_qdrant_upsert():
    chunks = [Chunk(chunk_id="a_000_000", text="hi", source_type="txt", source_name="a.txt")]
    vectors = [[0.1] * EMBEDDING_DIM]
    qdrant = _fake_qdrant_client()

    upsert_chunks(chunks, vectors, qdrant, "col")

    qdrant.upsert.assert_called_once()


def test_upsert_chunks_raises_on_length_mismatch():
    chunks = [Chunk(chunk_id="a_000_000", text="hi", source_type="txt", source_name="a.txt")]
    vectors = [[0.1] * EMBEDDING_DIM, [0.2] * EMBEDDING_DIM]
    qdrant = _fake_qdrant_client()

    with pytest.raises(ValueError, match="length mismatch"):
        upsert_chunks(chunks, vectors, qdrant, "col")


def test_upsert_chunks_empty_is_noop():
    qdrant = _fake_qdrant_client()
    upsert_chunks([], [], qdrant, "col")
    qdrant.upsert.assert_not_called()


def test_upsert_chunks_point_id_is_deterministic_uuid_not_random_hash():
    """Point IDs must not depend on Python's randomized hash() seed — otherwise
    re-processing the same chunk_id in a different process (restart, worker)
    produces a different point id and the old point is never overwritten,
    leaving stale duplicates in Qdrant instead of an idempotent upsert.

    uuid5 is deterministic by spec (RFC 4122) regardless of PYTHONHASHSEED,
    so pinning the point id to a uuid5 of chunk_id proves cross-process
    stability without needing to spawn a subprocess."""
    import uuid as _uuid

    chunk = Chunk(chunk_id="src1_000_000", text="hi", source_type="txt", source_name="a.txt")
    vectors = [[0.1] * EMBEDDING_DIM]
    qdrant = _fake_qdrant_client()

    upsert_chunks([chunk], vectors, qdrant, "col")

    point_id = qdrant.upsert.call_args.kwargs["points"][0].id
    # Must parse as a UUID (Qdrant requires int or UUID-string point ids).
    _uuid.UUID(str(point_id))
    # Must be a pure, deterministic function of chunk_id alone.
    upsert_chunks([chunk], vectors, (qdrant2 := _fake_qdrant_client()), "col")
    assert point_id == qdrant2.upsert.call_args.kwargs["points"][0].id


def test_upsert_chunks_different_chunk_ids_produce_different_point_ids():
    chunks = [
        Chunk(chunk_id="src1_000_000", text="hi", source_type="txt", source_name="a.txt"),
        Chunk(chunk_id="src2_000_000", text="yo", source_type="txt", source_name="a.txt"),
    ]
    vectors = [[0.1] * EMBEDDING_DIM, [0.2] * EMBEDDING_DIM]
    qdrant = _fake_qdrant_client()

    upsert_chunks(chunks, vectors, qdrant, "col")

    points = qdrant.upsert.call_args.kwargs["points"]
    assert points[0].id != points[1].id


# ---------------------------------------------------------------------------
# /upload route tests
# ---------------------------------------------------------------------------

def test_upload_pdf_returns_source_id_and_segment_count():
    pdf_bytes = _make_pdf_bytes("Hello from PDF page one.")
    with TestClient(app) as client:
        resp = client.post(
            "/upload",
            files={"file": ("doc.pdf", pdf_bytes, "application/pdf")},
        )
    assert resp.status_code == 200
    body = resp.json()
    assert body["source_type"] == "pdf"
    assert body["segment_count"] >= 1
    assert uuid.UUID(body["source_id"])  # valid UUID


def test_upload_txt_returns_source_id():
    with TestClient(app) as client:
        resp = client.post(
            "/upload",
            files={"file": ("notes.txt", b"Some text content here.", "text/plain")},
        )
    assert resp.status_code == 200
    assert resp.json()["source_type"] == "txt"


def test_upload_pptx_returns_source_id():
    pptx_bytes = _make_pptx_bytes("Slide heading text")
    with TestClient(app) as client:
        resp = client.post(
            "/upload",
            files={"file": ("deck.pptx", pptx_bytes, "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
        )
    assert resp.status_code == 200
    assert resp.json()["source_type"] == "pptx"


def test_upload_unsupported_extension_returns_422():
    with TestClient(app) as client:
        resp = client.post(
            "/upload",
            files={"file": ("data.csv", b"a,b,c", "text/csv")},
        )
    assert resp.status_code == 422
    assert "Unsupported" in resp.json()["detail"]


def test_upload_corrupted_pdf_returns_422():
    with TestClient(app) as client:
        resp = client.post(
            "/upload",
            files={"file": ("bad.pdf", b"not a pdf", "application/pdf")},
        )
    assert resp.status_code == 422


# ---------------------------------------------------------------------------
# /add-url route tests
# ---------------------------------------------------------------------------

def _url_client_override(html: str, status: int = 200):
    import httpx as _httpx

    async def _fake_get(url, **kwargs):
        return _httpx.Response(status, text=html)

    mock_client = MagicMock()
    mock_client.get = _fake_get
    return mock_client


def test_add_url_returns_source_id_for_valid_html():
    html = "<html><body><article>" + ("Real content sentence. " * 50) + "</article></body></html>"
    app.dependency_overrides[get_httpx_client] = lambda: _url_client_override(html)
    try:
        with TestClient(app) as client:
            resp = client.post("/add-url", json={"url": "https://example.com/article"})
        assert resp.status_code == 200
        body = resp.json()
        assert body["source_type"] == "url"
        assert uuid.UUID(body["source_id"])
    finally:
        app.dependency_overrides.pop(get_httpx_client, None)


def test_add_url_rejects_non_http_url():
    with TestClient(app) as client:
        resp = client.post("/add-url", json={"url": "ftp://example.com"})
    assert resp.status_code == 422


def test_add_url_returns_422_on_http_error():
    app.dependency_overrides[get_httpx_client] = lambda: _url_client_override("not found", status=404)
    try:
        with TestClient(app) as client:
            resp = client.post("/add-url", json={"url": "https://example.com/missing"})
        assert resp.status_code == 422
        assert "404" in resp.json()["detail"]
    finally:
        app.dependency_overrides.pop(get_httpx_client, None)


# ---------------------------------------------------------------------------
# /process route tests
# ---------------------------------------------------------------------------

def _setup_process_overrides(gemini=None, qdrant=None):
    app.dependency_overrides[get_gemini_client] = lambda: gemini or _fake_gemini_client()
    app.dependency_overrides[get_qdrant_client] = lambda: qdrant or _fake_qdrant_client()


def _teardown_process_overrides():
    app.dependency_overrides.pop(get_gemini_client, None)
    app.dependency_overrides.pop(get_qdrant_client, None)


def test_process_returns_job_id_immediately():
    # First upload a file to get a real source_id
    pdf_bytes = _make_pdf_bytes("Process test content.")
    _setup_process_overrides()
    try:
        with TestClient(app) as client:
            up = client.post("/upload", files={"file": ("p.pdf", pdf_bytes, "application/pdf")})
            source_id = up.json()["source_id"]

            resp = client.post("/process", json={"source_ids": [source_id]})
        assert resp.status_code == 200
        body = resp.json()
        assert "job_id" in body
        assert body["status"] in ("pending", "running", "done")
    finally:
        _teardown_process_overrides()


def test_process_unknown_source_id_returns_422():
    _setup_process_overrides()
    try:
        with TestClient(app) as client:
            resp = client.post("/process", json={"source_ids": ["nonexistent-id"]})
        assert resp.status_code == 422
        assert "Unknown source_ids" in resp.json()["detail"]
    finally:
        _teardown_process_overrides()


def test_process_empty_source_ids_returns_422():
    with TestClient(app) as client:
        resp = client.post("/process", json={"source_ids": []})
    assert resp.status_code == 422


def test_poll_process_returns_job_status():
    pdf_bytes = _make_pdf_bytes("Poll test content.")
    _setup_process_overrides()
    try:
        with TestClient(app) as client:
            up = client.post("/upload", files={"file": ("poll.pdf", pdf_bytes, "application/pdf")})
            source_id = up.json()["source_id"]

            proc = client.post("/process", json={"source_ids": [source_id]})
            job_id = proc.json()["job_id"]

            # Poll — background task may or may not have finished yet
            time.sleep(0.1)
            poll = client.get(f"/process/{job_id}")
        assert poll.status_code == 200
        body = poll.json()
        assert body["job_id"] == job_id
        assert body["status"] in ("pending", "running", "done", "error")
        assert len(body["sources"]) == 1
    finally:
        _teardown_process_overrides()


def test_poll_process_unknown_job_returns_404():
    with TestClient(app) as client:
        resp = client.get("/process/nonexistent-job-id")
    assert resp.status_code == 404


def test_process_completes_with_indexed_status():
    """End-to-end: upload → process → poll until done (mocked Gemini + Qdrant)."""
    pdf_bytes = _make_pdf_bytes("End to end test content for indexing.")
    _setup_process_overrides()
    try:
        with TestClient(app) as client:
            up = client.post("/upload", files={"file": ("e2e.pdf", pdf_bytes, "application/pdf")})
            source_id = up.json()["source_id"]

            proc = client.post("/process", json={"source_ids": [source_id]})
            job_id = proc.json()["job_id"]

            # Poll with retries — background task runs in the same thread with TestClient
            for _ in range(20):
                poll = client.get(f"/process/{job_id}")
                if poll.json()["status"] in ("done", "error"):
                    break
                time.sleep(0.05)

        body = poll.json()
        assert body["status"] == "done", f"Expected done, got: {body}"
        assert body["sources"][0]["status"] == "indexed"
    finally:
        _teardown_process_overrides()
