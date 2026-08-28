"""Local-file extraction service — turns an uploaded PDF/TXT/PPTX into
ExtractedSegments (raw text + source metadata), one segment per page,
slide, or whole file.

This is the first stage of the ingestion pipeline: upload -> **extract**
-> chunk (chunking.py) -> embed (embeddings.py) -> index (vector_store.py).
URL-based extraction is handled separately by app/services/scraper.py,
which produces the same ExtractedSegment shape so downstream code doesn't
need to know whether a source came from a file or a web page.
"""

from pathlib import Path

import fitz
from pptx import Presentation

from app.models import ExtractedSegment


class ExtractionError(Exception):
    """Raised when a source cannot be extracted. Message is human-readable."""


def extract_txt(file_path: Path, source_name: str) -> list[ExtractedSegment]:
    try:
        raw = file_path.read_text(encoding="utf-8")
    except UnicodeDecodeError as exc:
        raise ExtractionError(f"TXT file is not valid UTF-8: {source_name}") from exc

    text = "\n".join(line.strip() for line in raw.splitlines() if line.strip())

    if not text:
        raise ExtractionError(f"TXT file is empty: {source_name}")

    return [ExtractedSegment(text=text, source_type="txt", source_name=source_name)]


def extract_pdf(file_path: Path, source_name: str) -> list[ExtractedSegment]:
    try:
        doc = fitz.open(file_path)
    except Exception as exc:
        raise ExtractionError(f"Corrupted or unreadable PDF: {source_name} ({exc})") from exc

    segments: list[ExtractedSegment] = []
    try:
        for page_index in range(doc.page_count):
            page = doc.load_page(page_index)
            text = page.get_text().strip()
            if text:
                segments.append(ExtractedSegment(
                    text=text,
                    source_type="pdf",
                    source_name=source_name,
                    page=page_index + 1,
                ))
    finally:
        doc.close()

    if not segments:
        raise ExtractionError(f"PDF has no extractable text: {source_name}")

    return segments


def extract_pptx(file_path: Path, source_name: str) -> list[ExtractedSegment]:
    try:
        presentation = Presentation(file_path)
    except Exception as exc:
        raise ExtractionError(f"Corrupted or unreadable PPTX: {source_name} ({exc})") from exc

    segments: list[ExtractedSegment] = []
    for slide_index, slide in enumerate(presentation.slides):
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                shape_text = shape.text_frame.text.strip()
                if shape_text:
                    parts.append(shape_text)
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(
                        cell.text.strip() for cell in row.cells if cell.text.strip()
                    )
                    if row_text:
                        parts.append(row_text)

        slide_text = "\n".join(parts).strip()
        if slide_text:
            segments.append(ExtractedSegment(
                text=slide_text,
                source_type="pptx",
                source_name=source_name,
                slide=slide_index + 1,
            ))

    if not segments:
        raise ExtractionError(f"PPTX has no extractable text: {source_name}")

    return segments
